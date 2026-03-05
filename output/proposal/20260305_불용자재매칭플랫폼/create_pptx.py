#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
견적서 PPTX 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_proposal_pptx():
    # 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 색상 정의
    BLUE = RGBColor(30, 64, 175)  # Primary Blue
    GRAY = RGBColor(100, 100, 100)
    WHITE = RGBColor(255, 255, 255)

    # 슬라이드 1: 표지
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # 배경 박스
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()

    # 제목
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "개발 견적서"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 프로젝트명
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "불용자재 AI 매칭 플랫폼 MVP"
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 견적 금액
    txBox = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "개발 비용: 2,000만원 (부가세 별도)"
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 기간
    txBox = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "개발 기간: 8주 (2개월)"
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 날짜
    txBox = slide.shapes.add_textbox(Inches(7), Inches(7), Inches(2.5), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "2026년 3월 5일"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.RIGHT

    # 슬라이드 2: 견적 요약
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    title = slide.shapes.title
    title.text = "견적 요약"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE

    # 표 추가
    rows, cols = 8, 2
    left = Inches(1.5)
    top = Inches(2)
    width = Inches(7)
    height = Inches(4)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # 테이블 데이터
    data = [
        ["항목", "금액"],
        ["인건비 (3.35 M/M)", "1,900만원"],
        ["인프라/라이선스", "50만원"],
        ["기타 비용", "50만원"],
        ["합계 (부가세 별도)", "2,000만원"],
        ["부가세 (10%)", "200만원"],
        ["총액 (부가세 포함)", "2,200만원"],
        ["유지보수 (월, 선택)", "50만원"],
    ]

    # 테이블 스타일
    for i, row_data in enumerate(data):
        for j, text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = text

            # 헤더 스타일
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(18)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
            # 총액 강조
            elif i in [4, 6]:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(16)
                    paragraph.font.bold = True
            else:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(14)

    # 슬라이드 3: 개발 범위 (기능별)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "개발 범위 (8개 대분류)"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE

    # 표
    rows, cols = 9, 3
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.8), Inches(9), Inches(5)).table

    data = [
        ["구분", "주요 기능", "금액"],
        ["회원 관리", "회원가입(OCR), 로그인/인증", "300만원"],
        ["자재 관리", "엑셀 업로드, 목록, 상세", "640만원"],
        ["AI 매칭", "매칭 엔진, 결과 화면, 알림", "650만원"],
        ["거래 관리", "요청/승인, 이력", "270만원"],
        ["대시보드", "마이페이지 (통계)", "150만원"],
        ["UI/UX", "기본 디자인 커스터마이징", "175만원"],
        ["PM/인프라", "관리, 환경 구축, 배포", "330만원"],
        ["할인 후 합계", "", "2,000만원"],
    ]

    for i, row_data in enumerate(data):
        for j, text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = text

            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(16)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
            elif i == 8:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(14)
                    paragraph.font.bold = True
            else:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)

    # 슬라이드 4: 화면 구성
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "화면 구성 (7개 화면)"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE

    rows, cols = 8, 3
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.8), Inches(8), Inches(4.5)).table

    data = [
        ["#", "화면명", "주요 기능"],
        ["1", "회원가입/로그인", "OCR 인식, 이메일 로그인"],
        ["2", "자재 등록 (엑셀)", "엑셀 업로드, 템플릿, 검증"],
        ["3", "자재 목록", "내 등록 자재, 필터"],
        ["4", "자재 상세", "자재 정보, 거래 요청"],
        ["5", "AI 매칭 결과", "추천 자재, 매칭 점수"],
        ["6", "거래 관리", "요청/승인, 이력"],
        ["7", "마이페이지", "대시보드, 통계"],
    ]

    for i, row_data in enumerate(data):
        for j, text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = text

            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(16)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
            else:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(14)

    # 슬라이드 5: 개발 일정
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "개발 일정 (8주)"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE

    rows, cols = 9, 4
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.8), Inches(9), Inches(5)).table

    data = [
        ["주차", "단계", "주요 작업", "진척률"],
        ["1주", "분석/설계", "요구사항, DB 설계, 환경 구축", "12%"],
        ["2주", "회원", "회원가입(OCR), 로그인", "25%"],
        ["3주", "자재 등록", "엑셀 업로드, 자재 목록", "37%"],
        ["4주", "자재 조회", "자재 상세, 필터 [중간검수]", "50%"],
        ["5주", "AI 매칭 1", "매칭 엔진, 배치 작업", "62%"],
        ["6주", "AI 매칭 2", "매칭 결과, 알림 [AI검수]", "75%"],
        ["7주", "거래/대시보드", "거래 관리, 마이페이지", "87%"],
        ["8주", "테스트/배포", "테스트, 배포 [최종검수]", "100%"],
    ]

    for i, row_data in enumerate(data):
        for j, text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = text

            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(14)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
            elif "검수" in text:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 243, 205)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(11)
            else:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(11)

    # 슬라이드 6: 지급 조건
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "지급 조건"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE

    rows, cols = 6, 4
    table = slide.shapes.add_table(rows, cols, Inches(1.5), Inches(2), Inches(7), Inches(4)).table

    data = [
        ["구분", "시기", "비율", "금액 (부가세 별도)"],
        ["계약금", "계약 체결 시", "30%", "600만원"],
        ["1차 중도금", "4주차 (중간 검수)", "30%", "600만원"],
        ["2차 중도금", "6주차 (AI 매칭 완료)", "20%", "400만원"],
        ["잔금", "8주차 (최종 검수)", "20%", "400만원"],
        ["합계", "", "100%", "2,000만원"],
    ]

    for i, row_data in enumerate(data):
        for j, text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = text

            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(16)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
            elif i == 5:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(14)
                    paragraph.font.bold = True
            else:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(14)

    # 슬라이드 7: 포함/제외 범위
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "포함/제외 범위"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE

    # 포함 사항
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "✅ 포함 사항"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE

    included = [
        "회원가입/로그인 (OCR)",
        "엑셀 일괄 업로드",
        "AI 자동 매칭",
        "거래 관리",
        "마이페이지",
        "이메일 알림",
        "7개 화면 개발",
    ]

    for item in included:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.level = 0

    # 제외 사항
    txBox = slide.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(4), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "❌ 제외 사항 (Phase 2)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 50, 50)

    excluded = [
        "소셜 로그인 (+150만원)",
        "실시간 알림 (+200만원)",
        "고급 검색 (+200만원)",
        "실시간 채팅 (+400만원)",
        "이미지 인식 (+500만원)",
        "모바일 앱 (+1,500만원)",
    ]

    for item in excluded:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.level = 0

    # 슬라이드 8: 유지보수
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "유지보수 (선택 사항)"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE

    rows, cols = 5, 3
    table = slide.shapes.add_table(rows, cols, Inches(2), Inches(2), Inches(6), Inches(3)).table

    data = [
        ["항목", "시간", "비용"],
        ["인프라 운영", "-", "5만원"],
        ["기술 지원 (장애, 버그)", "월 3.5시간", "30만원"],
        ["기능 개선 (소규모)", "월 2시간", "15만원"],
        ["월간 합계", "월 5.5시간", "50만원"],
    ]

    for i, row_data in enumerate(data):
        for j, text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = text

            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(16)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
            elif i == 4:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(14)
                    paragraph.font.bold = True
            else:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(14)

    # 연간 정보
    txBox = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(6), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "연간 유지보수: 800만원 (부가세 별도)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "월간 50만원 × 12개월 + 분기 업데이트 200만원"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER

    # 슬라이드 9: 기술 스택
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "기술 스택"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE

    rows, cols = 7, 3
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(4.5)).table

    data = [
        ["구분", "기술", "선정 사유"],
        ["프론트엔드", "Next.js 14 + React 18", "SSR/SSG, SEO 최적화"],
        ["UI", "TailwindCSS + Shadcn/ui", "빠른 개발, 기본 템플릿"],
        ["백엔드", "Next.js API + Supabase", "서버리스, 통합 개발"],
        ["AI/ML", "OpenAI Embedding API", "텍스트 임베딩 (저비용)"],
        ["OCR", "Naver CLOVA OCR", "한국 사업자등록증 인식"],
        ["인프라", "Vercel + GitHub Actions", "자동 배포, CI/CD"],
    ]

    for i, row_data in enumerate(data):
        for j, text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = text

            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(16)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
            else:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(13)

    # 슬라이드 10: 마지막 (감사합니다)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()

    # 감사합니다
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "감사합니다"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 제안사 정보
    txBox = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(2))
    tf = txBox.text_frame

    info = [
        "제안사: [회사명]",
        "담당자: [담당자명]",
        "연락처: [전화번호]",
        "이메일: [이메일]",
    ]

    for i, line in enumerate(info):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # 저장
    output_path = "output/proposal/20260305_불용자재매칭플랫폼/견적서_개발범위기준.pptx"
    prs.save(output_path)
    print(f"✅ PowerPoint 파일 생성 완료: {output_path}")
    return output_path

if __name__ == "__main__":
    create_proposal_pptx()

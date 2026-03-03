#!/usr/bin/env python3
"""대건로지온 스마트 배차 관리 시스템 - 견적서 프레젠테이션 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import requests
from io import BytesIO

# ─── 색상 팔레트 ───
BLUE_DARK = RGBColor(0x0D, 0x27, 0x4D)    # #0D274D 진한 네이비
BLUE_PRIMARY = RGBColor(0x1E, 0x40, 0xAF)  # #1E40AF 메인 블루
BLUE_MEDIUM = RGBColor(0x25, 0x63, 0xEB)    # #2563EB (아래에서 재정의)
BLUE_LIGHT = RGBColor(0x3B, 0x82, 0xF6)    # #3B82F6 라이트 블루
BLUE_ACCENT = RGBColor(0x60, 0xA5, 0xFA)   # #60A5FA 액센트
BLUE_BG = RGBColor(0xEF, 0xF6, 0xFF)       # #EFF6FF 배경
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1F, 0x1F, 0x1F)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_BG = RGBColor(0xF3, 0xF4, 0xF6)
GREEN = RGBColor(0x05, 0x96, 0x69)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)

# Fix BLUE_MEDIUM
BLUE_MEDIUM = RGBColor(0x25, 0x63, 0xEB)   # #2563EB

# ─── Unsplash 이미지 다운로드 ───
IMAGE_URLS = {
    "logistics": "https://images.unsplash.com/photo-1586528116311-ad8dd3c8310d?w=800&q=80",
    "dashboard": "https://images.unsplash.com/photo-1551288049-bebda4e38f71?w=800&q=80",
    "problem": "https://images.unsplash.com/photo-1450101499163-c8848c66ca85?w=800&q=80",
    "solution": "https://images.unsplash.com/photo-1519389950473-47ba0277781c?w=800&q=80",
    "security": "https://images.unsplash.com/photo-1563986768609-322da13575f2?w=800&q=80",
    "ai": "https://images.unsplash.com/photo-1677442136019-21780ecad995?w=800&q=80",
    "api": "https://images.unsplash.com/photo-1558494949-ef010cbdcc31?w=800&q=80",
    "quote": "https://images.unsplash.com/photo-1554224155-6726b3ff858f?w=800&q=80",
    "timeline": "https://images.unsplash.com/photo-1506784983877-45594efa4cbe?w=800&q=80",
    "roi": "https://images.unsplash.com/photo-1460925895917-afdab827c52f?w=800&q=80",
    "team": "https://images.unsplash.com/photo-1522071820081-009f0129c71c?w=800&q=80",
    "thankyou": "https://images.unsplash.com/photo-1497366216548-37526070297c?w=800&q=80",
}

def download_image(url):
    """URL에서 이미지 다운로드"""
    try:
        resp = requests.get(url, timeout=10)
        if resp.status_code == 200:
            return BytesIO(resp.content)
    except:
        pass
    return None

def set_shape_alpha(shape, alpha_val):
    """도형 투명도 설정 (alpha_val: 0~100000, 100000=불투명)"""
    from pptx.oxml.ns import qn
    from lxml import etree
    spPr = shape._element.spPr
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is not None:
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is not None:
            for ex in srgbClr.findall(qn('a:alpha')):
                srgbClr.remove(ex)
            alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha_elem.set('val', str(alpha_val))

def add_bg_shape(slide, left, top, width, height, color, alpha=None):
    """배경 도형 추가"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        set_shape_alpha(shape, alpha)
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, color=BLACK,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_text(slide, left, top, width, height, items, font_size=12, color=BLACK):
    """불릿 텍스트 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "맑은 고딕"
        p.space_after = Pt(6)
        p.level = 0
    return txBox

def add_icon_circle(slide, left, top, size, color, text="", font_size=16):
    """원형 아이콘 추가"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "맑은 고딕"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
    return shape

def add_image_with_overlay(slide, img_stream, left, top, width, height, overlay_alpha=True):
    """이미지 + 오버레이 추가"""
    if img_stream:
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, left, top, width, height)
    if overlay_alpha:
        add_bg_shape(slide, left, top, width, height, BLUE_DARK, alpha=40000)


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    SW = prs.slide_width
    SH = prs.slide_height

    # 이미지 다운로드
    print("이미지 다운로드 중...")
    images = {}
    for key, url in IMAGE_URLS.items():
        print(f"  {key}...")
        images[key] = download_image(url)
    print("이미지 다운로드 완료")

    # ═══════════════════════════════════════════
    # SLIDE 1: 표지
    # ═══════════════════════════════════════════
    print("슬라이드 1: 표지")
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # 배경 이미지
    if images.get("logistics"):
        images["logistics"].seek(0)
        slide.shapes.add_picture(images["logistics"], Inches(0), Inches(0), SW, SH)

    # 어두운 오버레이
    overlay = add_bg_shape(slide, Inches(0), Inches(0), SW, SH, BLUE_DARK, alpha=55000)

    # 좌측 블루 바
    add_bg_shape(slide, Inches(0), Inches(0), Inches(0.15), SH, BLUE_LIGHT)

    # 상단 회사 로고 영역
    add_text_box(slide, Inches(1), Inches(0.8), Inches(5), Inches(0.5),
                 "[제안사명]", font_size=14, color=BLUE_ACCENT, bold=False)

    # 메인 타이틀
    add_text_box(slide, Inches(1), Inches(2.2), Inches(8), Inches(1.2),
                 "대건로지온", font_size=52, color=WHITE, bold=True)
    add_text_box(slide, Inches(1), Inches(3.4), Inches(10), Inches(0.8),
                 "스마트 배차 관리 시스템 구축 제안서", font_size=32, color=BLUE_ACCENT, bold=False)

    # 구분선
    add_bg_shape(slide, Inches(1), Inches(4.4), Inches(3), Inches(0.04), BLUE_LIGHT)

    # 하단 정보
    info_items = [
        "프로젝트 기간  |  3개월",
        "총 투자 금액  |  35,000,000원 (VAT 별도)",
        "제안일  |  2026년 3월 3일",
    ]
    for i, item in enumerate(info_items):
        add_text_box(slide, Inches(1), Inches(4.8 + i * 0.45), Inches(8), Inches(0.4),
                     item, font_size=16, color=RGBColor(0xD1, 0xD5, 0xDB))

    # ═══════════════════════════════════════════
    # SLIDE 2: 목차
    # ═══════════════════════════════════════════
    print("슬라이드 2: 목차")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 흰색 배경
    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)

    # 좌측 블루 사이드바
    add_bg_shape(slide, Inches(0), Inches(0), Inches(4.5), SH, BLUE_DARK)

    # 사이드바 텍스트
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(3), Inches(0.6),
                 "CONTENTS", font_size=14, color=BLUE_ACCENT, bold=True)
    add_text_box(slide, Inches(0.8), Inches(2.2), Inches(3), Inches(0.8),
                 "목차", font_size=36, color=WHITE, bold=True)
    add_bg_shape(slide, Inches(0.8), Inches(3.1), Inches(1.5), Inches(0.03), BLUE_LIGHT)

    # 목차 항목들
    toc_items = [
        ("01", "프로젝트 배경 및 목적"),
        ("02", "개발 범위 및 핵심 기능"),
        ("03", "사용자/권한/고객사 관리"),
        ("04", "AI 챗봇 & 플랫폼 통합"),
        ("05", "견적서 자동화 & 대시보드"),
        ("06", "시스템 아키텍처"),
        ("07", "개발 일정"),
        ("08", "투자 비용 및 기대 효과"),
        ("09", "협력 사항 및 결론"),
    ]

    for i, (num, title) in enumerate(toc_items):
        y = Inches(1.2 + i * 0.6)
        # 번호
        add_text_box(slide, Inches(5.3), y, Inches(0.8), Inches(0.5),
                     num, font_size=22, color=BLUE_PRIMARY, bold=True)
        # 구분 점
        add_icon_circle(slide, Inches(6.1), y + Inches(0.12), Inches(0.12), BLUE_LIGHT)
        # 제목
        add_text_box(slide, Inches(6.5), y, Inches(5.5), Inches(0.5),
                     title, font_size=16, color=BLACK)
        # 하단 라인
        if i < len(toc_items) - 1:
            add_bg_shape(slide, Inches(5.3), y + Inches(0.5), Inches(6.8), Inches(0.01), GRAY_BG)

    # ═══════════════════════════════════════════
    # SLIDE 3: 프로젝트 배경
    # ═══════════════════════════════════════════
    print("슬라이드 3: 프로젝트 배경")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)

    # 상단 헤더 바
    add_bg_shape(slide, Inches(0), Inches(0), SW, Inches(1.2), BLUE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(3), Inches(0.3),
                 "01", font_size=14, color=BLUE_ACCENT, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "프로젝트 배경 및 현재 상황", font_size=26, color=WHITE, bold=True)

    # 우측 이미지
    if images.get("problem"):
        images["problem"].seek(0)
        slide.shapes.add_picture(images["problem"], Inches(8.5), Inches(1.5), Inches(4.3), Inches(2.8))

    # 현재 상황 카드들
    problems = [
        ("다채널 접수", "카카오톡으로 견적 요청을\n개별 접수·관리"),
        ("수동 조회", "4개 화물 플랫폼을\n각각 개별 확인"),
        ("수작업 비교", "각 플랫폼 견적을\n수동으로 비교·선정"),
        ("분산 관리", "고객사 정보가\n엑셀·수첩에 산재"),
    ]

    for i, (title, desc) in enumerate(problems):
        x = Inches(0.6 + i * 1.9)
        y = Inches(1.6)
        card = add_bg_shape(slide, x, y, Inches(1.75), Inches(2.5), BLUE_BG)
        card.shadow.inherit = False
        add_icon_circle(slide, x + Inches(0.55), y + Inches(0.25), Inches(0.6), BLUE_PRIMARY, str(i+1), 18)
        add_text_box(slide, x + Inches(0.15), y + Inches(1.1), Inches(1.5), Inches(0.35),
                     title, font_size=13, color=BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), y + Inches(1.5), Inches(1.55), Inches(0.9),
                     desc, font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

    # 하단 해결 과제
    add_bg_shape(slide, Inches(0), Inches(4.6), SW, Inches(2.9), GRAY_BG)
    add_text_box(slide, Inches(0.8), Inches(4.8), Inches(5), Inches(0.4),
                 "핵심 해결 과제", font_size=18, color=BLUE_PRIMARY, bold=True)

    challenges = [
        "사용자 및 권한 관리 부재 — 직원별 역할과 권한이 체계적으로 관리되지 않음",
        "고객사 정보 분산 — 여러 곳에 흩어진 고객사 정보 관리 어려움",
        "분산된 업무 환경 — 카카오톡 + 4개 화물 플랫폼 각각 확인하는 비효율성",
        "수동 데이터 입력 — 동일 배차 정보를 4개 사이트에 반복 입력",
        "견적서 작성 시간 — 고객사별 수동 작성 반복 작업",
    ]
    for i, ch in enumerate(challenges):
        y = Inches(5.3 + i * 0.4)
        add_icon_circle(slide, Inches(0.9), y + Inches(0.05), Inches(0.18), ORANGE)
        add_text_box(slide, Inches(1.3), y, Inches(11), Inches(0.35),
                     ch, font_size=11, color=BLACK)

    # ═══════════════════════════════════════════
    # SLIDE 4: 프로젝트 목표
    # ═══════════════════════════════════════════
    print("슬라이드 4: 프로젝트 목표")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)

    # 헤더
    add_bg_shape(slide, Inches(0), Inches(0), SW, Inches(1.2), BLUE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(3), Inches(0.3),
                 "01", font_size=14, color=BLUE_ACCENT, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "프로젝트 목표 — 통합 배차 관리 플랫폼 구축", font_size=26, color=WHITE, bold=True)

    # 이미지
    if images.get("solution"):
        images["solution"].seek(0)
        pic = slide.shapes.add_picture(images["solution"], Inches(0), Inches(1.2), Inches(13.333), Inches(2.5))
        # 오버레이
        add_bg_shape(slide, Inches(0), Inches(1.2), SW, Inches(2.5), BLUE_DARK, alpha=60000)

    # 목표 아이콘 카드들 (이미지 위에)
    goals = [
        ("🔐", "계정/권한 관리", "체계적인 RBAC 기반\n사용자 권한 제어"),
        ("🏢", "고객사 CRM", "통합 고객사 관리로\n정보 일원화"),
        ("📂", "파일 자동화", "카카오톡 견적 파일\n자동 업로드·분류"),
        ("🤖", "AI 챗봇", "실시간 메뉴얼\n질의응답 시스템"),
        ("🔗", "플랫폼 통합", "4개 화물 플랫폼\n한 번에 조회"),
        ("📄", "견적서 자동화", "표준 견적서 자동 생성\n작성 시간 90% 단축"),
    ]

    for i, (icon, title, desc) in enumerate(goals):
        x = Inches(0.5 + i * 2.1)
        y = Inches(1.5)
        card = add_bg_shape(slide, x, y, Inches(1.9), Inches(1.9), WHITE)
        # 카드 내용
        add_text_box(slide, x + Inches(0.1), y + Inches(0.15), Inches(1.7), Inches(0.5),
                     icon, font_size=28, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.7), Inches(1.7), Inches(0.35),
                     title, font_size=12, color=BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.05), y + Inches(1.1), Inches(1.8), Inches(0.7),
                     desc, font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)

    # 하단 기대 효과 요약
    add_bg_shape(slide, Inches(0), Inches(4.0), SW, Inches(3.5), BLUE_BG)
    add_text_box(slide, Inches(0.8), Inches(4.2), Inches(5), Inches(0.4),
                 "핵심 기대 효과", font_size=18, color=BLUE_PRIMARY, bold=True)

    effects = [
        ("업무 시간", "70%", "절감"),
        ("처리량", "200%", "증가"),
        ("견적서 작성", "87%", "단축"),
        ("정보 검색", "95%", "단축"),
    ]
    for i, (label, value, unit) in enumerate(effects):
        x = Inches(0.8 + i * 3.1)
        y = Inches(4.8)
        add_bg_shape(slide, x, y, Inches(2.7), Inches(2.2), WHITE)
        add_text_box(slide, x, y + Inches(0.2), Inches(2.7), Inches(0.3),
                     label, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, y + Inches(0.6), Inches(2.7), Inches(0.8),
                     value, font_size=44, color=BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, y + Inches(1.5), Inches(2.7), Inches(0.3),
                     unit, font_size=14, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════
    # SLIDE 5: 개발 범위 및 핵심 기능
    # ═══════════════════════════════════════════
    print("슬라이드 5: 개발 범위")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)

    # 헤더
    add_bg_shape(slide, Inches(0), Inches(0), SW, Inches(1.2), BLUE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(3), Inches(0.3),
                 "02", font_size=14, color=BLUE_ACCENT, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "개발 범위 — 6대 핵심 시스템", font_size=26, color=WHITE, bold=True)

    # 6개 기능 카드
    features = [
        ("1", "사용자/권한/고객사\n관리 시스템", "계정 관리, RBAC 권한,\n고객사 CRM 통합", BLUE_PRIMARY, "⭐ 최우선"),
        ("2", "파일 업로드\n시스템", "카톡 파일 업로드,\n자동 분류, OCR", BLUE_MEDIUM, ""),
        ("3", "AI 챗봇\n메뉴얼 시스템", "GPT/Gemini 연동,\nRAG 기반 질의응답", BLUE_MEDIUM, ""),
        ("4", "4개 화물 플랫폼\nAPI 통합", "통합 입력/조회,\n실시간 견적 비교", BLUE_MEDIUM, ""),
        ("5", "통합 관리\n대시보드", "배차 현황, 통계,\n권한별 모니터링", BLUE_LIGHT, ""),
        ("6", "견적서 자동\n작성 시스템", "템플릿 기반 생성,\nPDF 출력, 이메일", BLUE_LIGHT, ""),
    ]

    for i, (num, title, desc, color, badge) in enumerate(features):
        col = i % 3
        row = i // 3
        x = Inches(0.6 + col * 4.1)
        y = Inches(1.5 + row * 3.0)

        # 카드 배경
        add_bg_shape(slide, x, y, Inches(3.8), Inches(2.7), GRAY_BG)

        # 상단 컬러 바
        add_bg_shape(slide, x, y, Inches(3.8), Inches(0.08), color)

        # 번호 원
        add_icon_circle(slide, x + Inches(0.2), y + Inches(0.3), Inches(0.55), color, num, 20)

        # 제목
        add_text_box(slide, x + Inches(0.9), y + Inches(0.3), Inches(2.7), Inches(0.7),
                     title, font_size=13, color=BLACK, bold=True)

        # 설명
        add_text_box(slide, x + Inches(0.2), y + Inches(1.3), Inches(3.4), Inches(1.0),
                     desc, font_size=11, color=GRAY)

        # 배지
        if badge:
            badge_shape = add_bg_shape(slide, x + Inches(2.5), y + Inches(0.15), Inches(1.1), Inches(0.3), ORANGE)
            add_text_box(slide, x + Inches(2.5), y + Inches(0.15), Inches(1.1), Inches(0.3),
                         badge, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════
    # SLIDE 6: 사용자/권한/고객사 관리
    # ═══════════════════════════════════════════
    print("슬라이드 6: 사용자/권한/고객사 관리")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)

    # 헤더
    add_bg_shape(slide, Inches(0), Inches(0), SW, Inches(1.2), BLUE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(3), Inches(0.3),
                 "03", font_size=14, color=BLUE_ACCENT, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "사용자 · 권한 · 고객사 관리 시스템", font_size=26, color=WHITE, bold=True)

    # 우측 이미지
    if images.get("security"):
        images["security"].seek(0)
        slide.shapes.add_picture(images["security"], Inches(8.8), Inches(1.4), Inches(4.2), Inches(2.8))

    # 좌측 3개 섹션
    sections = [
        ("계정 관리", [
            "이메일/아이디 기반 회원가입·로그인 (JWT)",
            "비밀번호 암호화 저장 (bcrypt)",
            "직원 계정 생성/수정/삭제, CSV 일괄 등록",
            "로그인 실패 제한 (5회), 세션 관리",
        ]),
        ("권한 관리 (RBAC)", [
            "슈퍼관리자 / 관리자 / 일반사용자 / 조회전용",
            "메뉴별·기능별(CRUD) 접근 권한 설정",
            "데이터 범위 권한 (전체/부서/본인)",
            "부서별 권한 관리",
        ]),
        ("고객사 CRM", [
            "기본정보·연락처·주소·담당자 통합 관리",
            "등급(VIP/A/B/C), 태그, 상태 분류",
            "계약/가격 정책, 거래 이력 추적",
            "통합 검색·필터링, 고객사 대시보드",
        ]),
    ]

    for i, (title, items) in enumerate(sections):
        y = Inches(1.5 + i * 1.9)
        add_bg_shape(slide, Inches(0.5), y, Inches(0.08), Inches(1.6), BLUE_PRIMARY)
        add_text_box(slide, Inches(0.8), y, Inches(4), Inches(0.35),
                     title, font_size=14, color=BLUE_PRIMARY, bold=True)
        for j, item in enumerate(items):
            add_text_box(slide, Inches(1.0), y + Inches(0.4 + j * 0.3), Inches(7.5), Inches(0.28),
                         f"•  {item}", font_size=10, color=BLACK)

    # 하단 권한 테이블 요약
    add_bg_shape(slide, Inches(0), Inches(7.0), SW, Inches(0.5), BLUE_BG)
    add_text_box(slide, Inches(0.8), Inches(7.05), Inches(12), Inches(0.4),
                 "감사 로그: 모든 로그인 이력, 데이터 변경 기록, 중요 작업 로그를 자동 기록 → 보안 감사 및 문제 추적",
                 font_size=10, color=BLUE_PRIMARY)

    # ═══════════════════════════════════════════
    # SLIDE 7: AI 챗봇 & 플랫폼 통합
    # ═══════════════════════════════════════════
    print("슬라이드 7: AI 챗봇 & 플랫폼 통합")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)

    # 헤더
    add_bg_shape(slide, Inches(0), Inches(0), SW, Inches(1.2), BLUE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(3), Inches(0.3),
                 "04", font_size=14, color=BLUE_ACCENT, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "AI 챗봇 메뉴얼 & 4개 화물 플랫폼 통합", font_size=26, color=WHITE, bold=True)

    # 좌측: AI 챗봇
    add_bg_shape(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5), BLUE_BG)
    add_text_box(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.4),
                 "🤖  AI 챗봇 메뉴얼 시스템", font_size=18, color=BLUE_PRIMARY, bold=True)

    if images.get("ai"):
        images["ai"].seek(0)
        slide.shapes.add_picture(images["ai"], Inches(3.8), Inches(2.3), Inches(2.3), Inches(1.5))

    ai_items = [
        "부서별 업무 메뉴얼 DB 구축 (벡터 DB)",
        "ChatGPT / Gemini API 멀티 LLM 지원",
        "RAG 기반 컨텍스트 답변 + 출처 표시",
        "자주 묻는 질문 Quick Button",
        "부서별 필터링, 질문 히스토리",
        "답변 유용성 피드백 수집",
    ]
    for i, item in enumerate(ai_items):
        add_text_box(slide, Inches(1.0), Inches(2.4 + i * 0.35), Inches(2.6), Inches(0.3),
                     f"•  {item}", font_size=10, color=BLACK)

    ai_stack = "LangChain · Pinecone · OpenAI · Gemini · FastAPI"
    add_text_box(slide, Inches(0.8), Inches(6.3), Inches(5), Inches(0.3),
                 ai_stack, font_size=9, color=GRAY)

    # 우측: 플랫폼 통합
    add_bg_shape(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(5.5), GRAY_BG)
    add_text_box(slide, Inches(7.1), Inches(1.7), Inches(5), Inches(0.4),
                 "🔗  4개 화물 플랫폼 API 통합", font_size=18, color=BLUE_PRIMARY, bold=True)

    if images.get("api"):
        images["api"].seek(0)
        slide.shapes.add_picture(images["api"], Inches(10.2), Inches(2.3), Inches(2.3), Inches(1.5))

    platforms = [
        "전국24시콜화물 · 화물맨 · 원콜 · 인성데이타",
        "",
        "한 번의 입력으로 4개 플랫폼 동시 조회",
        "실시간 견적 비교 대시보드",
        "가격 오름차순 정렬, 조건별 필터링",
        "기본 운임 + 추가 요금 상세 표시",
        "API / 웹 스크래핑 하이브리드 연동",
        "에러 핸들링 + 5분 캐싱 (Redis)",
    ]
    for i, item in enumerate(platforms):
        if item:
            add_text_box(slide, Inches(7.3), Inches(2.4 + i * 0.35), Inches(2.7), Inches(0.3),
                         f"•  {item}", font_size=10, color=BLACK)

    plat_stack = "Playwright · REST API · Redis · Node.js · Python"
    add_text_box(slide, Inches(7.1), Inches(6.3), Inches(5), Inches(0.3),
                 plat_stack, font_size=9, color=GRAY)

    # ═══════════════════════════════════════════
    # SLIDE 8: 견적서 자동화 & 대시보드
    # ═══════════════════════════════════════════
    print("슬라이드 8: 견적서 자동화 & 대시보드")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)

    # 헤더
    add_bg_shape(slide, Inches(0), Inches(0), SW, Inches(1.2), BLUE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(3), Inches(0.3),
                 "05", font_size=14, color=BLUE_ACCENT, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "견적서 자동 작성 시스템 & 통합 대시보드", font_size=26, color=WHITE, bold=True)

    # 좌측: 견적서 자동화
    add_bg_shape(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(5.5), WHITE)
    add_bg_shape(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(0.08), BLUE_PRIMARY)

    if images.get("quote"):
        images["quote"].seek(0)
        slide.shapes.add_picture(images["quote"], Inches(4), Inches(1.8), Inches(2.3), Inches(1.5))

    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(3), Inches(0.4),
                 "📄  견적서 자동 작성", font_size=18, color=BLUE_PRIMARY, bold=True)

    quote_items = [
        "템플릿 기반 자동 생성",
        "고객사 정보 CRM 자동 연동",
        "할인율/마진율 자동 계산",
        "PDF 고품질 출력",
        "이메일 자동 발송",
        "견적서 이력/승인/재발행 관리",
        "고객사별 커스터마이징",
    ]
    for i, item in enumerate(quote_items):
        add_text_box(slide, Inches(1.0), Inches(2.5 + i * 0.35), Inches(2.8), Inches(0.3),
                     f"•  {item}", font_size=10, color=BLACK)

    add_text_box(slide, Inches(0.8), Inches(5.2), Inches(5), Inches(0.4),
                 "견적서 작성 시간: 15분 → 2분 (87% 단축)",
                 font_size=12, color=GREEN, bold=True)

    # 우측: 대시보드
    add_bg_shape(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(5.5), BLUE_BG)
    add_bg_shape(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(0.08), BLUE_LIGHT)

    if images.get("dashboard"):
        images["dashboard"].seek(0)
        slide.shapes.add_picture(images["dashboard"], Inches(10.2), Inches(1.8), Inches(2.3), Inches(1.5))

    add_text_box(slide, Inches(7.3), Inches(1.8), Inches(3), Inches(0.4),
                 "📊  통합 대시보드", font_size=18, color=BLUE_PRIMARY, bold=True)

    dash_items = [
        "오늘의 배차 현황 (접수/처리중/완료)",
        "4개 플랫폼 이용 현황",
        "고객사 현황 (총 수, 신규)",
        "견적서 발행 현황 (일/주/월)",
        "AI 챗봇 인기 질문 TOP 5",
        "담당자별 업무 현황",
        "권한별 맞춤 위젯 표시",
    ]
    for i, item in enumerate(dash_items):
        add_text_box(slide, Inches(7.5), Inches(2.5 + i * 0.35), Inches(5), Inches(0.3),
                     f"•  {item}", font_size=10, color=BLACK)

    # 통계/리포트
    add_text_box(slide, Inches(7.3), Inches(5.2), Inches(5), Inches(0.35),
                 "통계/리포트", font_size=13, color=BLUE_PRIMARY, bold=True)
    stats = ["일별/주별/월별 처리 건수", "고객사별 거래·매출 분석", "Excel/CSV/PDF 데이터 내보내기"]
    for i, item in enumerate(stats):
        add_text_box(slide, Inches(7.5), Inches(5.6 + i * 0.3), Inches(5), Inches(0.25),
                     f"•  {item}", font_size=10, color=BLACK)

    # ═══════════════════════════════════════════
    # SLIDE 9: 시스템 아키텍처
    # ═══════════════════════════════════════════
    print("슬라이드 9: 시스템 아키텍처")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)

    # 헤더
    add_bg_shape(slide, Inches(0), Inches(0), SW, Inches(1.2), BLUE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(3), Inches(0.3),
                 "06", font_size=14, color=BLUE_ACCENT, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "시스템 아키텍처 & 기술 스택", font_size=26, color=WHITE, bold=True)

    # 아키텍처 다이어그램 (도형으로 구성)
    # 상단: 인증/권한
    add_bg_shape(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(0.7), BLUE_PRIMARY)
    add_text_box(slide, Inches(1.5), Inches(1.55), Inches(10.3), Inches(0.6),
                 "사용자 인증 및 권한 관리  (JWT · RBAC · bcrypt)", font_size=13, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # 중간: 핵심 모듈 (2x3 그리드)
    modules = [
        ("고객사 CRM", BLUE_MEDIUM), ("파일 관리", BLUE_MEDIUM), ("AI 챗봇", BLUE_MEDIUM),
        ("플랫폼 통합", BLUE_LIGHT), ("견적서 자동화", BLUE_LIGHT), ("통합 대시보드", BLUE_LIGHT),
    ]

    for i, (name, color) in enumerate(modules):
        col = i % 3
        row = i // 3
        x = Inches(1.5 + col * 3.5)
        y = Inches(2.5 + row * 1.2)
        add_bg_shape(slide, x, y, Inches(3.1), Inches(0.9), color)
        add_text_box(slide, x, y + Inches(0.1), Inches(3.1), Inches(0.7),
                     name, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 감사 로그
    add_bg_shape(slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.6), GRAY_BG)
    add_text_box(slide, Inches(1.5), Inches(5.05), Inches(10.3), Inches(0.5),
                 "감사 로그 시스템  (모든 활동 기록·추적)", font_size=12, color=GRAY, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # 하단: 기술 스택 카드
    add_bg_shape(slide, Inches(0), Inches(5.9), SW, Inches(1.6), BLUE_BG)
    add_text_box(slide, Inches(0.8), Inches(5.95), Inches(3), Inches(0.3),
                 "기술 스택", font_size=14, color=BLUE_PRIMARY, bold=True)

    stacks = [
        ("Frontend", "React 18 · TypeScript\nTailwind CSS · shadcn/ui\nChart.js · Zustand"),
        ("Backend", "Node.js · Express\nPython · FastAPI\nJWT · Passport.js"),
        ("Database", "PostgreSQL 14+\nPinecone (Vector DB)\nRedis 7+ (Cache)"),
        ("AI / ML", "OpenAI GPT-4\nGoogle Gemini\nLangChain · RAG"),
        ("Infra", "AWS / NCP\nDocker · S3\nPlaywright"),
    ]

    for i, (title, desc) in enumerate(stacks):
        x = Inches(0.6 + i * 2.55)
        add_bg_shape(slide, x, Inches(6.3), Inches(2.3), Inches(1.1), WHITE)
        add_text_box(slide, x, Inches(6.35), Inches(2.3), Inches(0.25),
                     title, font_size=10, color=BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), Inches(6.6), Inches(2.1), Inches(0.7),
                     desc, font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════
    # SLIDE 10: 개발 일정
    # ═══════════════════════════════════════════
    print("슬라이드 10: 개발 일정")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)

    # 헤더
    add_bg_shape(slide, Inches(0), Inches(0), SW, Inches(1.2), BLUE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(3), Inches(0.3),
                 "07", font_size=14, color=BLUE_ACCENT, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "개발 일정 — 3개월 타임라인", font_size=26, color=WHITE, bold=True)

    # 타임라인 이미지 배경
    if images.get("timeline"):
        images["timeline"].seek(0)
        slide.shapes.add_picture(images["timeline"], Inches(9.5), Inches(1.4), Inches(3.5), Inches(2.2))

    # 타임라인 바
    timeline_y = Inches(2.2)
    add_bg_shape(slide, Inches(0.8), timeline_y, Inches(8), Inches(0.06), BLUE_LIGHT)

    phases = [
        ("1-2주", "착수 · 분석", "킥오프, 요구사항 분석\n권한/역할 정의\n고객사 데이터 수집\n메뉴얼 데이터 수집", BLUE_PRIMARY),
        ("3-4주", "환경 구축", "개발 서버 세팅\nDB 스키마 설계\nAPI 구조 설계\nUI/UX 디자인", BLUE_MEDIUM),
        ("5-7주", "1단계 개발", "사용자/권한/고객사\n파일 업로드 시스템\nAI 챗봇 시스템\n벡터 DB 구축", BLUE_LIGHT),
        ("8-10주", "2단계 개발", "4개 플랫폼 API 연동\n통합 대시보드\n견적서 자동 작성\n견적 비교 기능", BLUE_LIGHT),
        ("11주", "통합 테스트", "통합 테스트\n권한별 테스트\n성능 최적화\n보안 점검", GREEN),
        ("12주", "배포 · 교육", "사용자 교육\n운영 서버 배포\n매뉴얼 전달\n프로젝트 완료", GREEN),
    ]

    for i, (week, title, desc, color) in enumerate(phases):
        x = Inches(0.5 + i * 1.4)
        # 타임라인 점
        add_icon_circle(slide, x + Inches(0.45), timeline_y - Inches(0.08), Inches(0.2), color)

        # 주차
        add_text_box(slide, x, Inches(1.7), Inches(1.3), Inches(0.25),
                     week, font_size=10, color=color, bold=True, alignment=PP_ALIGN.CENTER)

        # 단계명 카드
        add_bg_shape(slide, x, Inches(2.6), Inches(1.3), Inches(0.4), color)
        add_text_box(slide, x, Inches(2.62), Inches(1.3), Inches(0.35),
                     title, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # 상세 내용
        add_text_box(slide, x, Inches(3.1), Inches(1.3), Inches(1.5),
                     desc, font_size=8, color=GRAY, alignment=PP_ALIGN.CENTER)

    # 마일스톤 요약
    add_bg_shape(slide, Inches(0), Inches(4.8), SW, Inches(2.7), BLUE_BG)
    add_text_box(slide, Inches(0.8), Inches(5.0), Inches(5), Inches(0.4),
                 "주요 마일스톤", font_size=16, color=BLUE_PRIMARY, bold=True)

    milestones = [
        ("M1  2주차", "요구사항 확정, 권한/역할 정의, 견적서 템플릿 승인"),
        ("M2  4주차", "개발 환경 구축 완료, DB 스키마 확정"),
        ("M3  7주차", "사용자/권한/고객사 관리 + 파일 업로드 + AI 챗봇 완료"),
        ("M4  10주차", "4개 플랫폼 통합 + 견적서 자동 작성 완료"),
        ("M5  12주차", "통합 테스트 완료 → 운영 배포 → 인수인계"),
    ]
    for i, (ms, desc) in enumerate(milestones):
        y = Inches(5.5 + i * 0.38)
        add_bg_shape(slide, Inches(0.9), y, Inches(1.2), Inches(0.3), BLUE_PRIMARY)
        add_text_box(slide, Inches(0.9), y, Inches(1.2), Inches(0.3),
                     ms, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(2.3), y, Inches(10), Inches(0.3),
                     desc, font_size=10, color=BLACK)

    # ═══════════════════════════════════════════
    # SLIDE 11: 투자 비용
    # ═══════════════════════════════════════════
    print("슬라이드 11: 투자 비용")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)

    # 헤더
    add_bg_shape(slide, Inches(0), Inches(0), SW, Inches(1.2), BLUE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(3), Inches(0.3),
                 "08", font_size=14, color=BLUE_ACCENT, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "투자 비용", font_size=26, color=WHITE, bold=True)

    # 메인 금액 카드
    add_bg_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5), BLUE_PRIMARY)
    add_text_box(slide, Inches(1.2), Inches(1.7), Inches(4.5), Inches(0.4),
                 "총 프로젝트 비용", font_size=16, color=BLUE_ACCENT)
    add_text_box(slide, Inches(1.2), Inches(2.2), Inches(4.5), Inches(0.8),
                 "35,000,000원", font_size=42, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.2), Inches(3.1), Inches(4.5), Inches(0.4),
                 "VAT 별도  |  3개월 개발", font_size=14, color=RGBColor(0xBF, 0xDB, 0xFE))

    # 우측 이미지
    if images.get("roi"):
        images["roi"].seek(0)
        slide.shapes.add_picture(images["roi"], Inches(6.8), Inches(1.5), Inches(6), Inches(2.5))

    # 비용 내역
    add_text_box(slide, Inches(0.8), Inches(4.3), Inches(5), Inches(0.4),
                 "비용 내역", font_size=16, color=BLUE_PRIMARY, bold=True)

    cost_items = [
        ("인건비", "PM 1명 + FE 1명 + BE 2명 + AI 1명 + QA 1명", "28,000,000원"),
        ("인프라/라이선스", "클라우드, AI API, 개발 도구", "2,500,000원"),
        ("프로젝트 관리비", "형상관리, 문서화, 리스크 관리", "2,500,000원"),
        ("QA/테스트 비용", "통합 테스트, 보안 점검", "1,000,000원"),
        ("예비비", "요구사항 변경 대응", "1,000,000원"),
    ]

    for i, (item, desc, amount) in enumerate(cost_items):
        y = Inches(4.8 + i * 0.5)
        add_bg_shape(slide, Inches(0.8), y, Inches(5.5), Inches(0.42), GRAY_BG if i % 2 == 0 else WHITE)
        add_text_box(slide, Inches(1.0), y + Inches(0.02), Inches(1.5), Inches(0.35),
                     item, font_size=11, color=BLACK, bold=True)
        add_text_box(slide, Inches(2.5), y + Inches(0.02), Inches(2.5), Inches(0.35),
                     desc, font_size=9, color=GRAY)
        add_text_box(slide, Inches(4.8), y + Inches(0.02), Inches(1.5), Inches(0.35),
                     amount, font_size=11, color=BLACK, bold=True, alignment=PP_ALIGN.RIGHT)

    # 우측: 기대 효과
    add_bg_shape(slide, Inches(7), Inches(4.3), Inches(5.8), Inches(3), BLUE_BG)
    add_text_box(slide, Inches(7.3), Inches(4.5), Inches(5), Inches(0.4),
                 "투자 대비 기대 효과 (연간)", font_size=16, color=BLUE_PRIMARY, bold=True)

    roi_items = [
        ("인건비 절감", "업무 시간 60% 절감", "약 3,500만원/년"),
        ("매출 증대", "처리량 200% 증가", "약 8,000만원/년"),
        ("비용 절감", "견적서 외주 + 고객 이탈 방지", "약 1,500만원/년"),
    ]
    for i, (title, desc, amount) in enumerate(roi_items):
        y = Inches(5.0 + i * 0.65)
        add_bg_shape(slide, Inches(7.3), y, Inches(5.2), Inches(0.55), WHITE)
        add_text_box(slide, Inches(7.5), y + Inches(0.02), Inches(1.3), Inches(0.25),
                     title, font_size=11, color=BLUE_PRIMARY, bold=True)
        add_text_box(slide, Inches(7.5), y + Inches(0.28), Inches(2), Inches(0.2),
                     desc, font_size=9, color=GRAY)
        add_text_box(slide, Inches(10.5), y + Inches(0.1), Inches(1.8), Inches(0.3),
                     amount, font_size=13, color=GREEN, bold=True, alignment=PP_ALIGN.RIGHT)

    add_text_box(slide, Inches(7.3), Inches(7.0), Inches(5.2), Inches(0.3),
                 "투자 회수 기간: 약 3~4개월", font_size=13, color=BLUE_PRIMARY, bold=True)

    # ═══════════════════════════════════════════
    # SLIDE 12: 협력 사항
    # ═══════════════════════════════════════════
    print("슬라이드 12: 협력 사항")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)

    # 헤더
    add_bg_shape(slide, Inches(0), Inches(0), SW, Inches(1.2), BLUE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(3), Inches(0.3),
                 "09", font_size=14, color=BLUE_ACCENT, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "협력 사항 및 프로젝트 결론", font_size=26, color=WHITE, bold=True)

    if images.get("team"):
        images["team"].seek(0)
        slide.shapes.add_picture(images["team"], Inches(8.5), Inches(1.5), Inches(4.5), Inches(2.5))

    # 대건로지온 준비 사항
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.4),
                 "대건로지온 준비 사항", font_size=16, color=BLUE_PRIMARY, bold=True)

    client_items = [
        "조직/권한 정보 (부서 구조, 직원 목록, 권한 요구)",
        "고객사 데이터 (Excel/CSV), 분류 기준, 가격 정책",
        "부서별 업무 메뉴얼, 요일별 업무 흐름, FAQ",
        "4개 화물 플랫폼 계정 정보, API 키 발급 협조",
        "기존 견적서 샘플 (10건+), 브랜딩 자료 (로고, 컬러)",
        "프로젝트 담당자 지정, 주 1회 정기 미팅 참석",
    ]
    for i, item in enumerate(client_items):
        add_text_box(slide, Inches(1.0), Inches(2.0 + i * 0.35), Inches(7), Inches(0.3),
                     f"•  {item}", font_size=10, color=BLACK)

    # 제안사 제공
    add_text_box(slide, Inches(0.8), Inches(4.3), Inches(5), Inches(0.4),
                 "제안사 제공 사항", font_size=16, color=BLUE_PRIMARY, bold=True)

    vendor_items = [
        "PM 1명 (풀타임) + 개발팀 5명",
        "주간 진행 상황 보고, 이슈/리스크 관리",
        "소스 코드 (GitHub), API 문서, DB 스키마 문서",
        "운영 매뉴얼, 사용자 가이드 (권한별)",
        "3개월 무상 유지보수 (배포 후)",
    ]
    for i, item in enumerate(vendor_items):
        add_text_box(slide, Inches(1.0), Inches(4.8 + i * 0.35), Inches(7), Inches(0.3),
                     f"•  {item}", font_size=10, color=BLACK)

    # 우측 하단: 유지보수
    add_bg_shape(slide, Inches(8), Inches(4.3), Inches(4.8), Inches(3), BLUE_BG)
    add_text_box(slide, Inches(8.3), Inches(4.5), Inches(4), Inches(0.35),
                 "운영 및 유지보수", font_size=14, color=BLUE_PRIMARY, bold=True)

    maint_items = [
        ("무상 (3개월)", "버그 수정, 장애 대응\n시스템 모니터링, 기술 지원"),
        ("유상 (4개월~)", "월 50~100만원 (협의)\n보안 패치, 기능 개선\n기술 지원 월 20시간"),
    ]
    for i, (title, desc) in enumerate(maint_items):
        y = Inches(5.0 + i * 1.2)
        add_bg_shape(slide, Inches(8.3), y, Inches(4.2), Inches(1.0), WHITE)
        add_text_box(slide, Inches(8.5), y + Inches(0.05), Inches(1.5), Inches(0.25),
                     title, font_size=11, color=BLUE_PRIMARY, bold=True)
        add_text_box(slide, Inches(10.0), y + Inches(0.05), Inches(2.3), Inches(0.8),
                     desc, font_size=9, color=GRAY)

    # ═══════════════════════════════════════════
    # SLIDE 13: Thank You
    # ═══════════════════════════════════════════
    print("슬라이드 13: Thank You")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경 이미지
    if images.get("thankyou"):
        images["thankyou"].seek(0)
        slide.shapes.add_picture(images["thankyou"], Inches(0), Inches(0), SW, SH)

    # 어두운 오버레이
    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, BLUE_DARK, alpha=60000)

    # 좌측 블루 바
    add_bg_shape(slide, Inches(0), Inches(0), Inches(0.15), SH, BLUE_LIGHT)

    add_text_box(slide, Inches(0), Inches(1.8), SW, Inches(0.8),
                 "Thank You", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0), Inches(2.8), SW, Inches(0.5),
                 "대건로지온 스마트 배차 관리 시스템 구축 제안", font_size=20, color=BLUE_ACCENT,
                 alignment=PP_ALIGN.CENTER)

    # 구분선
    add_bg_shape(slide, Inches(5.5), Inches(3.6), Inches(2.3), Inches(0.03), BLUE_LIGHT)

    # 연락처
    add_text_box(slide, Inches(0), Inches(4.2), SW, Inches(0.4),
                 "[제안사명]", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    contact_lines = [
        "담당자: [이름]",
        "이메일: [email@example.com]",
        "전화: [010-XXXX-XXXX]",
    ]
    for i, line in enumerate(contact_lines):
        add_text_box(slide, Inches(0), Inches(4.8 + i * 0.4), SW, Inches(0.35),
                     line, font_size=14, color=RGBColor(0xD1, 0xD5, 0xDB), alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0), Inches(6.2), SW, Inches(0.3),
                 "2026년 3월 3일", font_size=12, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)

    # ─── 저장 ───
    output_path = "/Users/jasonmac/kitt_agent/ai-agent/output/proposal/대건로지온_견적서_프레젠테이션.pptx"
    prs.save(output_path)
    print(f"\n✅ 프레젠테이션 생성 완료: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()

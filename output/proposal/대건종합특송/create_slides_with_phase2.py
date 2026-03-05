#!/usr/bin/env python3
"""대건종합특송 스마트 배차 관리 시스템 - 견적서 프레젠테이션 생성 (2차 개발 포함)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import requests
from io import BytesIO

# ─── 색상 팔레트 (화이트 + 블루 전문가 스타일) ───
# Primary Colors
NAVY_DARK = RGBColor(0x0D, 0x27, 0x4D)      # #0D274D
BLUE_PRIMARY = RGBColor(0x1E, 0x40, 0xAF)   # #1E40AF
BLUE_MEDIUM = RGBColor(0x3B, 0x82, 0xF6)    # #3B82F6
BLUE_LIGHT = RGBColor(0x60, 0xA5, 0xFA)     # #60A5FA
BLUE_ACCENT = RGBColor(0x93, 0xC5, 0xFD)    # #93C5FD

# Background Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)          # #FFFFFF
GRAY_50 = RGBColor(0xF9, 0xFA, 0xFB)        # #F9FAFB
BLUE_50 = RGBColor(0xEF, 0xF6, 0xFF)        # #EFF6FF
GRAY_100 = RGBColor(0xF3, 0xF4, 0xF6)       # #F3F4F6

# Semantic Colors
SUCCESS = RGBColor(0x10, 0xB9, 0x81)        # #10B981
WARNING = RGBColor(0xF5, 0x9E, 0x0B)        # #F59E0B
DANGER = RGBColor(0xDC, 0x26, 0x26)         # #DC2626

# Text Colors
TEXT_DARK = RGBColor(0x0D, 0x27, 0x4D)      # #0D274D
TEXT_PRIMARY = RGBColor(0x37, 0x41, 0x51)   # #374151
TEXT_SECONDARY = RGBColor(0x6B, 0x72, 0x80) # #6B7280
TEXT_LIGHT = RGBColor(0x9C, 0xA3, 0xAF)     # #9CA3AF

# Border Colors
BORDER = RGBColor(0xE5, 0xE7, 0xEB)         # #E5E7EB
BORDER_BLUE = RGBColor(0x3B, 0x82, 0xF6)    # #3B82F6

# Legacy aliases
BLUE_DARK = NAVY_DARK
BLUE_BG = BLUE_50
BLACK = TEXT_DARK
GRAY = TEXT_SECONDARY
GRAY_LIGHT = TEXT_LIGHT
GRAY_BG = GRAY_50
GREEN = SUCCESS
ORANGE = WARNING

# ─── 슬라이드 크기 ───
SW = Inches(13.33)  # 16:9
SH = Inches(7.5)

# ─── Unsplash 이미지 다운로드 ───
IMAGE_URLS = {
    "logistics": "https://images.unsplash.com/photo-1586528116311-ad8dd3c8310d?w=800&q=80",
    "dashboard": "https://images.unsplash.com/photo-1551288049-bebda4e38f71?w=800&q=80",
    "problem": "https://images.unsplash.com/photo-1450101499163-c8848c66ca85?w=800&q=80",
    "solution": "https://images.unsplash.com/photo-1519389950473-47ba0277781c?w=800&q=80",
    "security": "https://images.unsplash.com/photo-1563986768609-322da13575f2?w=800&q=80",
    "ai": "https://images.unsplash.com/photo-1677442136019-21780ecad995?w=800&q=80",
    "api": "https://images.unsplash.com/photo-1558494949-ef010cbdcc31?w=800&q=80",
    "quote": "https://images.unsplash.com/photo-1554224155-6726b3ff858f?w=800&q=80",
    "timeline": "https://images.unsplash.com/photo-1506784983877-45594efa4cbe?w=800&q=80",
    "roi": "https://images.unsplash.com/photo-1460925895917-afdab827c52f?w=800&q=80",
    "team": "https://images.unsplash.com/photo-1522071820081-009f0129c71c?w=800&q=80",
    "mobile": "https://images.unsplash.com/photo-1512941937669-90a1b58e7e9c?w=800&q=80",
    "tracking": "https://images.unsplash.com/photo-1586880244406-556ebe35f282?w=800&q=80",
}

def download_image(url):
    """URL에서 이미지 다운로드"""
    try:
        resp = requests.get(url, timeout=10)
        if resp.status_code == 200:
            return BytesIO(resp.content)
    except:
        pass
    return None

def set_shape_alpha(shape, alpha_val):
    """도형 투명도 설정 (alpha_val: 0~100000, 100000=불투명)"""
    from pptx.oxml.ns import qn
    from lxml import etree
    spPr = shape._element.spPr
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is not None:
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is not None:
            for ex in srgbClr.findall(qn('a:alpha')):
                srgbClr.remove(ex)
            alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha_elem.set('val', str(alpha_val))

def add_bg_shape(slide, left, top, width, height, color, alpha=None):
    """배경 도형 추가"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        set_shape_alpha(shape, alpha)
    return shape

def add_left_blue_line(slide):
    """좌측 블루 라인 추가 (4px)"""
    return add_bg_shape(slide, Inches(0), Inches(0), Inches(0.04), SH, BLUE_PRIMARY)

def add_text_box(slide, left, top, width, height, text, font_size=14, color=TEXT_DARK,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_text(slide, left, top, width, height, items, font_size=12, color=TEXT_PRIMARY):
    """불릿 텍스트 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "맑은 고딕"
        p.level = 0
    return txBox

def create_presentation():
    """프레젠테이션 생성"""
    print("프레젠테이션 생성 시작...")

    # 이미지 다운로드
    print("\n이미지 다운로드 중...")
    images = {}
    for key, url in IMAGE_URLS.items():
        print(f"  - {key}: {url}")
        img = download_image(url)
        if img:
            images[key] = img
            print(f"    ✓ 다운로드 완료")
        else:
            print(f"    ✗ 실패")

    # 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    # ═══════════════════════════════════════════
    # SLIDE 1: 표지 (화이트 배경)
    # ═══════════════════════════════════════════
    print("\n슬라이드 01: 표지")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경: 순백색
    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)

    # 좌측 블루 라인 (4px)
    add_left_blue_line(slide)

    # 상단 태그
    add_text_box(slide, Inches(0.8), Inches(0.8), Inches(5), Inches(0.3),
                 "SYSTEM PROPOSAL", font_size=12, color=BLUE_LIGHT, bold=True)

    # 메인 타이틀
    add_text_box(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(0.6),
                 "대건종합특송", font_size=48, color=NAVY_DARK, bold=True)

    add_text_box(slide, Inches(0.8), Inches(2.9), Inches(11), Inches(0.5),
                 "스마트 배차 관리 시스템 구축 제안서", font_size=28, color=BLUE_PRIMARY)

    # 구분선
    add_bg_shape(slide, Inches(0.8), Inches(3.6), Inches(3), Inches(0.02), BORDER_BLUE)

    # 부제
    add_text_box(slide, Inches(0.8), Inches(4.0), Inches(8), Inches(0.4),
                 "🚚 대형차 배차 업무 디지털 전환", font_size=18, color=TEXT_SECONDARY)

    # 하단 정보
    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(5), Inches(0.3),
                 "[제안사명]", font_size=16, color=TEXT_PRIMARY, bold=True)
    add_text_box(slide, Inches(0.8), Inches(6.9), Inches(5), Inches(0.3),
                 "2026년 3월 4일", font_size=14, color=TEXT_LIGHT)

    # ═══════════════════════════════════════════
    # SLIDE 2: 1차 개발 범위 (기존 내용)
    # ═══════════════════════════════════════════
    print("슬라이드 02: 1차 개발 범위")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)
    add_left_blue_line(slide)

    # 타이틀
    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.4),
                 "📦 1차 개발 범위 (3개월, 35,000,000원)", font_size=24, color=NAVY_DARK, bold=True)

    # 왼쪽 열
    items_left = [
        ("1️⃣ 사용자 및 고객사 관리", "3,000,000원", [
            "계정 관리 (로그인/회원가입/권한)",
            "RBAC 권한 시스템",
            "고객사 CRM (정보/거래/이력)",
            "감사 로그 (활동 추적)"
        ]),
        ("2️⃣ 파일 업로드 시스템", "5,000,000원", [
            "카톡 파일 업로드",
            "자동 분류 및 상태 관리",
            "OCR 데이터 추출 (선택)"
        ]),
        ("3️⃣ AI 챗봇 메뉴얼 시스템", "8,500,000원", [
            "벡터 DB 구축 (RAG)",
            "ChatGPT/Gemini API 연동",
            "실시간 질의응답"
        ])
    ]

    y_pos = 1.2
    for title, price, details in items_left:
        # 박스 높이 동적 계산
        box_height = 0.5 + (len(details) * 0.22)

        # 박스 배경
        add_bg_shape(slide, Inches(0.8), Inches(y_pos), Inches(5.5), Inches(box_height), GRAY_50)

        # 타이틀 + 가격
        add_text_box(slide, Inches(1.0), Inches(y_pos + 0.08), Inches(3.8), Inches(0.25),
                     title, font_size=13, color=BLUE_PRIMARY, bold=True)
        add_text_box(slide, Inches(4.8), Inches(y_pos + 0.08), Inches(1.3), Inches(0.25),
                     price, font_size=12, color=SUCCESS, bold=True, alignment=PP_ALIGN.RIGHT)

        # 구분선
        add_bg_shape(slide, Inches(1.0), Inches(y_pos + 0.36), Inches(4.8), Inches(0.01), BORDER)

        # 세부 항목
        for i, detail in enumerate(details):
            add_text_box(slide, Inches(1.1), Inches(y_pos + 0.42 + i * 0.22), Inches(4.5), Inches(0.18),
                         f"• {detail}", font_size=9, color=TEXT_SECONDARY)

        y_pos += box_height + 0.15

    # 오른쪽 열
    items_right = [
        ("4️⃣ 4개 플랫폼 API 통합", "10,000,000원", [
            "전국24시콜화물 • 화물맨",
            "원콜 • 인성데이타",
            "통합 입력 + 실시간 비교"
        ]),
        ("5️⃣ 통합 관리 대시보드", "3,500,000원", [
            "권한별 맞춤 대시보드",
            "실시간 모니터링",
            "통계 리포트"
        ]),
        ("6️⃣ 표준 견적서 자동 작성", "3,000,000원", [
            "고객사 정보 자동 연동 (CRM)",
            "할인율/마진율 자동 계산",
            "PDF 생성 + 이메일 발송"
        ]),
        ("7️⃣ 인프라 및 배포", "2,000,000원", [
            "AWS/NCP 서버 구축",
            "CI/CD 파이프라인",
            "보안 설정"
        ])
    ]

    y_pos = 1.2
    for title, price, details in items_right:
        # 박스 높이 동적 계산
        box_height = 0.5 + (len(details) * 0.22)

        # 박스 배경
        add_bg_shape(slide, Inches(6.8), Inches(y_pos), Inches(5.5), Inches(box_height), GRAY_50)

        # 타이틀 + 가격
        add_text_box(slide, Inches(7.0), Inches(y_pos + 0.08), Inches(3.8), Inches(0.25),
                     title, font_size=13, color=BLUE_PRIMARY, bold=True)
        add_text_box(slide, Inches(10.8), Inches(y_pos + 0.08), Inches(1.3), Inches(0.25),
                     price, font_size=12, color=SUCCESS, bold=True, alignment=PP_ALIGN.RIGHT)

        # 구분선
        add_bg_shape(slide, Inches(7.0), Inches(y_pos + 0.36), Inches(4.8), Inches(0.01), BORDER)

        # 세부 항목
        for i, detail in enumerate(details):
            add_text_box(slide, Inches(7.1), Inches(y_pos + 0.42 + i * 0.22), Inches(4.5), Inches(0.18),
                         f"• {detail}", font_size=9, color=TEXT_SECONDARY)

        y_pos += box_height + 0.15

    # 총액
    add_bg_shape(slide, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.55), BLUE_50)
    add_text_box(slide, Inches(1.0), Inches(6.65), Inches(9), Inches(0.4),
                 "총 개발비 (VAT 별도)", font_size=15, color=TEXT_DARK, bold=True)
    add_text_box(slide, Inches(10), Inches(6.65), Inches(2), Inches(0.4),
                 "35,000,000원", font_size=18, color=BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.RIGHT)

    # ═══════════════════════════════════════════
    # SLIDE 3: 2차 개발 범위 (선택사항) - 새로 추가
    # ═══════════════════════════════════════════
    print("슬라이드 03: 2차 개발 범위 (선택사항)")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)
    add_left_blue_line(slide)

    # 타이틀
    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.4),
                 "🚀 2차 개발 범위 (선택사항)", font_size=24, color=NAVY_DARK, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.3),
                 "1차 개발 완료 후 추가로 구축 가능한 고급 기능", font_size=13, color=TEXT_SECONDARY)

    # 왼쪽 열
    phase2_left = [
        ("1️⃣ 외주 차량 관리 시스템", "5,000,000원", [
            "외주 차량 등록/관리",
            "외주 업체 정보 관리",
            "정산 시스템 (수수료 계산)",
            "외주 차량 배차 이력"
        ]),
        ("2️⃣ 배송 추적 시스템", "7,000,000원", [
            "실시간 GPS 위치 추적",
            "배송 상태 관리 (픽업/배송중/완료)",
            "고객 SMS 자동 알림",
            "배송 이력 조회"
        ]),
        ("3️⃣ 고급 템플릿 관리", "4,000,000원", [
            "엑셀 템플릿 라이브러리",
            "자동 양식 생성 시스템",
            "템플릿 버전 관리",
            "사용자 정의 템플릿 등록"
        ])
    ]

    y_pos = 1.5
    for title, price, details in phase2_left:
        # 박스 높이 동적 계산
        box_height = 0.5 + (len(details) * 0.22)

        # 박스 배경 (연한 블루)
        add_bg_shape(slide, Inches(0.8), Inches(y_pos), Inches(5.5), Inches(box_height), BLUE_50)

        # 좌측 강조선 (블루)
        add_bg_shape(slide, Inches(0.8), Inches(y_pos), Inches(0.04), Inches(box_height), BORDER_BLUE)

        # 타이틀 + 가격
        add_text_box(slide, Inches(1.0), Inches(y_pos + 0.08), Inches(3.8), Inches(0.25),
                     title, font_size=13, color=BLUE_PRIMARY, bold=True)
        add_text_box(slide, Inches(4.8), Inches(y_pos + 0.08), Inches(1.3), Inches(0.25),
                     price, font_size=12, color=SUCCESS, bold=True, alignment=PP_ALIGN.RIGHT)

        # 구분선
        add_bg_shape(slide, Inches(1.0), Inches(y_pos + 0.36), Inches(4.8), Inches(0.01), BORDER_BLUE)

        # 세부 항목
        for i, detail in enumerate(details):
            add_text_box(slide, Inches(1.1), Inches(y_pos + 0.42 + i * 0.22), Inches(4.5), Inches(0.18),
                         f"• {detail}", font_size=9, color=TEXT_SECONDARY)

        y_pos += box_height + 0.15

    # 오른쪽 열
    phase2_right = [
        ("4️⃣ 모바일 앱 (기사용)", "8,000,000원", [
            "기사 전용 모바일 앱",
            "배차 알림 푸시",
            "GPS 위치 자동 전송",
            "전자 서명 기능"
        ]),
        ("5️⃣ 고급 CRM 기능", "6,000,000원", [
            "고객 분석 대시보드",
            "자동화된 마케팅 (이메일/SMS)",
            "맞춤형 리포트 생성",
            "고객 세그먼트 관리"
        ])
    ]

    y_pos = 1.5
    for title, price, details in phase2_right:
        # 박스 높이 동적 계산
        box_height = 0.5 + (len(details) * 0.22)

        # 박스 배경
        add_bg_shape(slide, Inches(6.8), Inches(y_pos), Inches(5.5), Inches(box_height), BLUE_50)

        # 좌측 강조선
        add_bg_shape(slide, Inches(6.8), Inches(y_pos), Inches(0.04), Inches(box_height), BORDER_BLUE)

        # 타이틀 + 가격
        add_text_box(slide, Inches(7.0), Inches(y_pos + 0.08), Inches(3.8), Inches(0.25),
                     title, font_size=13, color=BLUE_PRIMARY, bold=True)
        add_text_box(slide, Inches(10.8), Inches(y_pos + 0.08), Inches(1.3), Inches(0.25),
                     price, font_size=12, color=SUCCESS, bold=True, alignment=PP_ALIGN.RIGHT)

        # 구분선
        add_bg_shape(slide, Inches(7.0), Inches(y_pos + 0.36), Inches(4.8), Inches(0.01), BORDER_BLUE)

        # 세부 항목
        for i, detail in enumerate(details):
            add_text_box(slide, Inches(7.1), Inches(y_pos + 0.42 + i * 0.22), Inches(4.5), Inches(0.18),
                         f"• {detail}", font_size=9, color=TEXT_SECONDARY)

        y_pos += box_height + 0.15

    # 총액 (2차 개발)
    add_bg_shape(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.4), GRAY_50)
    add_text_box(slide, Inches(1.0), Inches(6.35), Inches(9), Inches(0.3),
                 "2차 개발 총액 (VAT 별도)", font_size=14, color=TEXT_DARK, bold=True)
    add_text_box(slide, Inches(10), Inches(6.35), Inches(2), Inches(0.3),
                 "30,000,000원", font_size=17, color=BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.RIGHT)

    # 안내 문구
    add_text_box(slide, Inches(0.8), Inches(6.85), Inches(11.5), Inches(0.25),
                 "※ 2차 개발은 1차 개발 완료 후 별도 협의를 통해 진행 가능합니다. (전체 또는 부분 선택 가능)",
                 font_size=9, color=WARNING)

    # ═══════════════════════════════════════════
    # SLIDE 4: Thank You
    # ═══════════════════════════════════════════
    print("슬라이드 04: Thank You")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경: 화이트
    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)

    # 좌측 블루 라인
    add_left_blue_line(slide)

    # Thank You
    add_text_box(slide, Inches(0), Inches(2.5), SW, Inches(0.8),
                 "Thank You", font_size=52, color=BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0), Inches(3.4), SW, Inches(0.4),
                 "대건종합특송 스마트 배차 관리 시스템 구축 제안", font_size=18, color=TEXT_SECONDARY,
                 alignment=PP_ALIGN.CENTER)

    # 구분선
    add_bg_shape(slide, Inches(5.5), Inches(4.0), Inches(2.3), Inches(0.02), BORDER_BLUE)

    # 연락처
    add_text_box(slide, Inches(0), Inches(4.5), SW, Inches(0.3),
                 "[제안사명]", font_size=18, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)

    contact_lines = [
        "담당자: [이름]",
        "이메일: [email@example.com]",
        "전화: [010-XXXX-XXXX]",
    ]
    for i, line in enumerate(contact_lines):
        add_text_box(slide, Inches(0), Inches(5.0 + i * 0.35), SW, Inches(0.3),
                     line, font_size=14, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0), Inches(6.5), SW, Inches(0.3),
                 "2026년 3월 4일", font_size=12, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

    # ─── 저장 ───
    output_path = "/Users/jasonmac/kitt_agent/ai-agent/output/proposal/대건종합특송_견적서_화이트블루.pptx"
    prs.save(output_path)
    print(f"\n✅ 프레젠테이션 생성 완료: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()

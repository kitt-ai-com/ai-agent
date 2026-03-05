#!/usr/bin/env python3
"""정동병원 재고관리 & 세금계산서 자동화 플랫폼 - 제안서 프레젠테이션 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── 색상 팔레트 (에메랄드 그린 전문가 스타일) ───
# Primary Colors
EMERALD_DARK = RGBColor(0x06, 0x54, 0x40)     # #065440
EMERALD_PRIMARY = RGBColor(0x10, 0xB9, 0x81)  # #10B981
EMERALD_MEDIUM = RGBColor(0x34, 0xD3, 0x99)   # #34D399
EMERALD_LIGHT = RGBColor(0x6E, 0xE7, 0xB7)    # #6EE7B7
EMERALD_ACCENT = RGBColor(0xA7, 0xF3, 0xD0)   # #A7F3D0

# Background Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)            # #FFFFFF
GRAY_50 = RGBColor(0xF9, 0xFA, 0xFB)          # #F9FAFB
EMERALD_50 = RGBColor(0xEC, 0xFD, 0xF5)       # #ECFDF5
GRAY_100 = RGBColor(0xF3, 0xF4, 0xF6)         # #F3F4F6

# Semantic Colors
SUCCESS = RGBColor(0x10, 0xB9, 0x81)          # #10B981
WARNING = RGBColor(0xF5, 0x9E, 0x0B)          # #F59E0B
BLUE = RGBColor(0x3B, 0x82, 0xF6)             # #3B82F6
PURPLE = RGBColor(0xA8, 0x55, 0xF7)           # #A855F7
ORANGE = RGBColor(0xF9, 0x73, 0x16)           # #F97316

# Text Colors
TEXT_DARK = RGBColor(0x0D, 0x27, 0x4D)        # #0D274D
TEXT_PRIMARY = RGBColor(0x37, 0x41, 0x51)     # #374151
TEXT_SECONDARY = RGBColor(0x6B, 0x72, 0x80)   # #6B7280
TEXT_LIGHT = RGBColor(0x9C, 0xA3, 0xAF)       # #9CA3AF

# Border Colors
BORDER = RGBColor(0xE5, 0xE7, 0xEB)           # #E5E7EB
BORDER_EMERALD = RGBColor(0x10, 0xB9, 0x81)   # #10B981

# ─── 슬라이드 크기 ───
SW = Inches(13.33)  # 16:9
SH = Inches(7.5)

def add_bg_shape(slide, left, top, width, height, color):
    """배경 도형 추가"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_left_emerald_line(slide):
    """좌측 에메랄드 라인 추가 (4px)"""
    return add_bg_shape(slide, Inches(0), Inches(0), Inches(0.04), SH, EMERALD_PRIMARY)

def add_text_box(slide, left, top, width, height, text, font_size=14, color=TEXT_DARK,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def create_presentation():
    """프레젠테이션 생성"""
    print("정동병원 제안서 프레젠테이션 생성 시작...")

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    # ═══════════════════════════════════════════
    # SLIDE 1: 표지
    # ═══════════════════════════════════════════
    print("\n슬라이드 01: 표지")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)
    add_left_emerald_line(slide)

    # 상단 태그
    add_text_box(slide, Inches(0.8), Inches(0.8), Inches(5), Inches(0.3),
                 "INVENTORY MANAGEMENT PLATFORM", font_size=11, color=EMERALD_LIGHT, bold=True)

    # 메인 타이틀
    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(0.6),
                 "정동병원", font_size=48, color=EMERALD_DARK, bold=True)

    add_text_box(slide, Inches(0.8), Inches(2.7), Inches(11), Inches(0.5),
                 "재고관리 & 세금계산서 자동화 플랫폼", font_size=32, color=EMERALD_PRIMARY)

    # 구분선
    add_bg_shape(slide, Inches(0.8), Inches(3.4), Inches(3.5), Inches(0.03), BORDER_EMERALD)

    # 부제
    add_text_box(slide, Inches(0.8), Inches(3.9), Inches(8), Inches(0.4),
                 "📦 병원 운영 효율화를 위한 스마트 솔루션", font_size=20, color=TEXT_SECONDARY)

    # 주요 수치
    metrics_y = 4.8
    metrics = [
        ("재고 부족 사고", "90% 예방", SUCCESS),
        ("세금계산서 발행", "75% 단축", BLUE),
        ("투자 회수 기간", "약 1년", PURPLE)
    ]

    for i, (label, value, color) in enumerate(metrics):
        x_pos = 0.8 + (i * 3.5)
        add_text_box(slide, Inches(x_pos), Inches(metrics_y), Inches(3), Inches(0.25),
                     label, font_size=11, color=TEXT_SECONDARY)
        add_text_box(slide, Inches(x_pos), Inches(metrics_y + 0.3), Inches(3), Inches(0.4),
                     value, font_size=22, color=color, bold=True)

    # 하단 정보
    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(5), Inches(0.3),
                 "kitt AI Intelligence", font_size=16, color=TEXT_PRIMARY, bold=True)
    add_text_box(slide, Inches(0.8), Inches(6.9), Inches(5), Inches(0.3),
                 "2026년 3월", font_size=14, color=TEXT_LIGHT)

    # ═══════════════════════════════════════════
    # SLIDE 2: 정동병원 맞춤 분석
    # ═══════════════════════════════════════════
    print("슬라이드 02: 정동병원 맞춤 분석")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)
    add_left_emerald_line(slide)

    # 타이틀
    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.4),
                 "🏥 정동병원 맞춤 제안", font_size=28, color=EMERALD_DARK, bold=True)

    # 좌측: 현황 분석
    add_bg_shape(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(3.5), EMERALD_50)
    add_text_box(slide, Inches(1.2), Inches(1.5), Inches(5), Inches(0.35),
                 "예상 현황 분석", font_size=18, color=EMERALD_PRIMARY, bold=True)

    hospital_info = [
        ("🛏️ 병상 규모", "100~200병상 (중소 규모)"),
        ("👥 연간 환자 수", "약 50,000명 추정"),
        ("🏥 진료과", "내과, 외과, 정형외과 등"),
        ("💡 디지털화 수준", "AI 음성봇 도입 검토 중")
    ]

    y_pos = 2.1
    for icon_title, detail in hospital_info:
        add_text_box(slide, Inches(1.3), Inches(y_pos), Inches(5), Inches(0.22),
                     icon_title, font_size=12, color=TEXT_PRIMARY, bold=True)
        add_text_box(slide, Inches(1.3), Inches(y_pos + 0.25), Inches(5), Inches(0.2),
                     detail, font_size=10, color=TEXT_SECONDARY)
        y_pos += 0.65

    # 우측 상단: 추천 플랜
    add_bg_shape(slide, Inches(7.0), Inches(1.3), Inches(5.3), Inches(1.8), EMERALD_PRIMARY)
    add_text_box(slide, Inches(7.4), Inches(1.6), Inches(4.5), Inches(0.3),
                 "⭐ 추천 플랜", font_size=16, color=WHITE, bold=True)
    add_text_box(slide, Inches(7.4), Inches(2.0), Inches(4.5), Inches(0.5),
                 "프로페셔널 (Pro)", font_size=28, color=WHITE, bold=True)
    add_text_box(slide, Inches(7.4), Inches(2.6), Inches(4.5), Inches(0.25),
                 "중소 규모 병원에 최적화된 플랜", font_size=11, color=EMERALD_LIGHT)

    # 우측 하단: ROI
    add_bg_shape(slide, Inches(7.0), Inches(3.3), Inches(5.3), Inches(1.5), GRAY_50)
    add_text_box(slide, Inches(7.4), Inches(3.5), Inches(4.5), Inches(0.3),
                 "💰 투자 수익 분석 (ROI)", font_size=14, color=TEXT_DARK, bold=True)

    roi_data = [
        ("1년차 총 투자", "15.7M", TEXT_PRIMARY),
        ("연간 절감 효과", "30M", SUCCESS),
        ("순 이익 (ROI)", "91%", PURPLE)
    ]

    x_start = 7.4
    for i, (label, value, color) in enumerate(roi_data):
        x_pos = x_start + (i * 1.65)
        add_text_box(slide, Inches(x_pos), Inches(4.0), Inches(1.5), Inches(0.2),
                     label, font_size=8, color=TEXT_SECONDARY)
        add_text_box(slide, Inches(x_pos), Inches(4.2), Inches(1.5), Inches(0.4),
                     value, font_size=18, color=color, bold=True)

    # 하단: 기대 효과
    add_bg_shape(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(2.0), GRAY_50)
    add_text_box(slide, Inches(1.2), Inches(5.3), Inches(10.7), Inches(0.3),
                 "📊 정동병원 예상 도입 효과", font_size=16, color=TEXT_DARK, bold=True)

    effects = [
        ("업무 효율 75% 향상", "하루 3시간 → 45분\n연간 약 900시간 절약"),
        ("재고 비용 30% 절감", "연간 약 3,000만원 절감\n약 1년 만에 투자 회수"),
        ("세금계산서 80% 개선", "월말 2일 → 2시간 완료\n정산 업무 부담 감소")
    ]

    x_start = 1.2
    for i, (title, detail) in enumerate(effects):
        x_pos = x_start + (i * 3.7)
        add_text_box(slide, Inches(x_pos), Inches(5.8), Inches(3.3), Inches(0.25),
                     title, font_size=12, color=EMERALD_PRIMARY, bold=True)
        add_text_box(slide, Inches(x_pos), Inches(6.1), Inches(3.3), Inches(0.8),
                     detail, font_size=9, color=TEXT_SECONDARY)

    # ═══════════════════════════════════════════
    # SLIDE 3: 4가지 핵심 지표
    # ═══════════════════════════════════════════
    print("슬라이드 03: 4가지 핵심 지표")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)
    add_left_emerald_line(slide)

    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.4),
                 "📈 도입 후 달라지는 4가지 핵심 지표", font_size=28, color=EMERALD_DARK, bold=True)

    indicators = [
        {
            "icon": "📦",
            "title": "실시간 재고 관리",
            "desc": "입고/출고 데이터 자동 동기화\n재고 파악 시간 90% 단축",
            "color": EMERALD_PRIMARY,
            "bg": EMERALD_50
        },
        {
            "icon": "💳",
            "title": "세금계산서 자동화",
            "desc": "자동 발행 및 국세청 연동\n원스톱 신고 처리",
            "color": BLUE,
            "bg": RGBColor(0xDC, 0xED, 0xFF)
        },
        {
            "icon": "🔔",
            "title": "스마트 알림 시스템",
            "desc": "AI 기반 재고 부족 예측\n실시간 발주 알림",
            "color": ORANGE,
            "bg": RGBColor(0xFF, 0xED, 0xD5)
        },
        {
            "icon": "📊",
            "title": "데이터 기반 의사결정",
            "desc": "소비 패턴 분석\n최적 발주량 및 예산 계획",
            "color": PURPLE,
            "bg": RGBColor(0xF3, 0xE8, 0xFF)
        }
    ]

    # 2x2 그리드
    positions = [
        (0.8, 1.4), (6.9, 1.4),
        (0.8, 4.3), (6.9, 4.3)
    ]

    for i, (indicator, pos) in enumerate(zip(indicators, positions)):
        x, y = pos
        box_w = 5.6
        box_h = 2.5

        # 배경
        add_bg_shape(slide, Inches(x), Inches(y), Inches(box_w), Inches(box_h), indicator["bg"])

        # 좌측 강조선
        add_bg_shape(slide, Inches(x), Inches(y), Inches(0.05), Inches(box_h), indicator["color"])

        # 아이콘
        add_text_box(slide, Inches(x + 0.3), Inches(y + 0.3), Inches(0.6), Inches(0.6),
                     indicator["icon"], font_size=32, color=indicator["color"])

        # 타이틀
        add_text_box(slide, Inches(x + 1.0), Inches(y + 0.35), Inches(4.3), Inches(0.4),
                     indicator["title"], font_size=16, color=indicator["color"], bold=True)

        # 설명
        add_text_box(slide, Inches(x + 0.3), Inches(y + 1.1), Inches(5.0), Inches(1.1),
                     indicator["desc"], font_size=11, color=TEXT_SECONDARY)

    # ═══════════════════════════════════════════
    # SLIDE 4: 플랫폼 핵심 기능
    # ═══════════════════════════════════════════
    print("슬라이드 04: 플랫폼 핵심 기능")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)
    add_left_emerald_line(slide)

    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.4),
                 "🔧 플랫폼 핵심 기능", font_size=28, color=EMERALD_DARK, bold=True)

    features = [
        {
            "title": "1️⃣ 통합 재고 관리",
            "items": [
                "바코드/QR 스캔 입출고",
                "유통기한 자동 관리",
                "부서별 분산 재고 통합",
                "모바일 앱 실시간 확인"
            ],
            "color": EMERALD_PRIMARY
        },
        {
            "title": "2️⃣ 세금계산서 자동화",
            "items": [
                "입고 데이터 기반 자동 생성",
                "홈택스 전자세금계산서 발행",
                "월별/분기별 정산 리포트",
                "회계 프로그램 연동"
            ],
            "color": BLUE
        },
        {
            "title": "3️⃣ 공급업체 협업 포털",
            "items": [
                "발주서 자동 전송",
                "배송 상태 실시간 추적",
                "가격 비교 및 협상 이력",
                "공급업체 평가 시스템"
            ],
            "color": ORANGE
        },
        {
            "title": "4️⃣ AI 분석 & 예측",
            "items": [
                "최적 발주량 AI 추천",
                "소비 트렌드 분석",
                "재고 회전율 모니터링",
                "이상 거래 자동 탐지"
            ],
            "color": PURPLE
        }
    ]

    # 2x2 그리드
    positions = [
        (0.8, 1.4), (6.9, 1.4),
        (0.8, 4.3), (6.9, 4.3)
    ]

    for feature, pos in zip(features, positions):
        x, y = pos
        box_w = 5.6
        box_h = 2.5

        # 배경
        add_bg_shape(slide, Inches(x), Inches(y), Inches(box_w), Inches(box_h), GRAY_50)

        # 타이틀 배경
        add_bg_shape(slide, Inches(x), Inches(y), Inches(box_w), Inches(0.5), feature["color"])

        # 타이틀
        add_text_box(slide, Inches(x + 0.2), Inches(y + 0.1), Inches(5.2), Inches(0.35),
                     feature["title"], font_size=14, color=WHITE, bold=True)

        # 항목들
        for i, item in enumerate(feature["items"]):
            add_text_box(slide, Inches(x + 0.25), Inches(y + 0.7 + i * 0.4), Inches(5.1), Inches(0.3),
                         f"✓ {item}", font_size=10, color=TEXT_SECONDARY)

    # ═══════════════════════════════════════════
    # SLIDE 5: 도입 프로세스
    # ═══════════════════════════════════════════
    print("슬라이드 05: 도입 프로세스")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)
    add_left_emerald_line(slide)

    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.4),
                 "⏱️ 정동병원 맞춤 도입 일정 (4주 완성)", font_size=28, color=EMERALD_DARK, bold=True)

    # 예상 오픈일
    add_bg_shape(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.5), EMERALD_50)
    add_text_box(slide, Inches(1.0), Inches(1.3), Inches(11.1), Inches(0.3),
                 "예상 오픈일: 2026년 4월 30일", font_size=16, color=EMERALD_PRIMARY, bold=True,
                 alignment=PP_ALIGN.CENTER)

    weeks = [
        {
            "week": "WEEK 01",
            "date": "4월 1주",
            "title": "현황 분석 &\n요구사항 정의",
            "desc": "프로세스 분석\n품목 분류 체계\n사용자 권한 설계",
            "color": EMERALD_DARK
        },
        {
            "week": "WEEK 02",
            "date": "4월 2주",
            "title": "시스템 구축 &\n데이터 이관",
            "desc": "재고 데이터 입력\n공급업체 등록\n바코드 부착",
            "color": EMERALD_PRIMARY
        },
        {
            "week": "WEEK 03",
            "date": "4월 3주",
            "title": "사용자 교육 &\n테스트",
            "desc": "담당자별 교육\n시범 운영\n피드백 반영",
            "color": EMERALD_PRIMARY
        },
        {
            "week": "WEEK 04",
            "date": "4월 4주",
            "title": "정식 오픈 &\n안정화",
            "desc": "시스템 가동\n실시간 모니터링\n최적화",
            "color": EMERALD_DARK
        }
    ]

    x_start = 0.8
    for i, week in enumerate(weeks):
        x = x_start + (i * 2.95)
        box_w = 2.8
        box_h = 4.5

        # 배경
        add_bg_shape(slide, Inches(x), Inches(2.0), Inches(box_w), Inches(box_h), GRAY_50)

        # 상단 라벨
        add_bg_shape(slide, Inches(x), Inches(2.0), Inches(box_w), Inches(0.5), week["color"])
        add_text_box(slide, Inches(x), Inches(2.05), Inches(box_w), Inches(0.4),
                     week["week"], font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # 날짜
        add_text_box(slide, Inches(x), Inches(2.6), Inches(box_w), Inches(0.3),
                     week["date"], font_size=11, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

        # 타이틀
        add_text_box(slide, Inches(x + 0.15), Inches(3.0), Inches(box_w - 0.3), Inches(0.7),
                     week["title"], font_size=13, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)

        # 구분선
        add_bg_shape(slide, Inches(x + 0.3), Inches(3.8), Inches(box_w - 0.6), Inches(0.015), BORDER)

        # 설명
        add_text_box(slide, Inches(x + 0.2), Inches(4.0), Inches(box_w - 0.4), Inches(2.2),
                     week["desc"], font_size=10, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════
    # SLIDE 6: 비용 안내
    # ═══════════════════════════════════════════
    print("슬라이드 06: 비용 안내")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)
    add_left_emerald_line(slide)

    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.4),
                 "💰 도입 비용 안내", font_size=28, color=EMERALD_DARK, bold=True)

    plans = [
        {
            "name": "스타터",
            "subtitle": "소규모 의원",
            "price": "4,500,000원",
            "monthly": "300,000원/월",
            "features": [
                "기본 재고 관리",
                "세금계산서 자동 발행",
                "모바일 앱",
                "최대 500개 품목",
                "최대 5명 사용자"
            ],
            "bg": GRAY_50,
            "color": TEXT_PRIMARY
        },
        {
            "name": "프로페셔널 ⭐",
            "subtitle": "중소 규모 병원 추천",
            "price": "8,500,000원",
            "monthly": "600,000원/월",
            "features": [
                "스타터 전체 기능",
                "공급업체 협업 포털",
                "AI 발주량 예측",
                "ERP/회계 연동",
                "무제한 품목/20명"
            ],
            "bg": EMERALD_PRIMARY,
            "color": WHITE
        },
        {
            "name": "엔터프라이즈",
            "subtitle": "대형 종합병원",
            "price": "별도 협의",
            "monthly": "별도 협의",
            "features": [
                "프로 전체 기능",
                "다지점 통합 관리",
                "맞춤형 워크플로우",
                "HIS 완전 통합",
                "전담 매니저"
            ],
            "bg": GRAY_50,
            "color": TEXT_PRIMARY
        }
    ]

    x_start = 0.8
    for i, plan in enumerate(plans):
        x = x_start + (i * 3.95)
        box_w = 3.7
        box_h = 5.5

        # 배경
        add_bg_shape(slide, Inches(x), Inches(1.3), Inches(box_w), Inches(box_h), plan["bg"])

        # 플랜명
        add_text_box(slide, Inches(x + 0.2), Inches(1.5), Inches(box_w - 0.4), Inches(0.35),
                     plan["name"], font_size=16, color=plan["color"], bold=True, alignment=PP_ALIGN.CENTER)

        # 부제
        add_text_box(slide, Inches(x + 0.2), Inches(1.9), Inches(box_w - 0.4), Inches(0.25),
                     plan["subtitle"], font_size=10, color=plan["color"], alignment=PP_ALIGN.CENTER)

        # 구분선
        line_color = WHITE if i == 1 else BORDER
        add_bg_shape(slide, Inches(x + 0.5), Inches(2.25), Inches(box_w - 1.0), Inches(0.015), line_color)

        # 가격
        add_text_box(slide, Inches(x + 0.2), Inches(2.4), Inches(box_w - 0.4), Inches(0.45),
                     plan["price"], font_size=20, color=plan["color"], bold=True, alignment=PP_ALIGN.CENTER)

        # 월 운영료
        add_text_box(slide, Inches(x + 0.2), Inches(2.9), Inches(box_w - 0.4), Inches(0.2),
                     f"월 {plan['monthly']}", font_size=10, color=plan["color"], alignment=PP_ALIGN.CENTER)

        # 기능 목록
        y_feature = 3.3
        for feature in plan["features"]:
            add_text_box(slide, Inches(x + 0.3), Inches(y_feature), Inches(box_w - 0.6), Inches(0.25),
                         f"✓ {feature}", font_size=9, color=plan["color"])
            y_feature += 0.35

    # ═══════════════════════════════════════════
    # SLIDE 7: 기대 효과
    # ═══════════════════════════════════════════
    print("슬라이드 07: 기대 효과")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)
    add_left_emerald_line(slide)

    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.4),
                 "✨ 정동병원 기대 효과", font_size=28, color=EMERALD_DARK, bold=True)

    # 좌측: 정량적 효과
    add_bg_shape(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(5.5), EMERALD_50)
    add_text_box(slide, Inches(1.2), Inches(1.6), Inches(5), Inches(0.35),
                 "📊 정량적 효과", font_size=18, color=EMERALD_PRIMARY, bold=True)

    quantitative = [
        ("📈 업무 효율 75% 향상", "하루 3시간 → 45분\n연간 약 900시간 절약", SUCCESS),
        ("💰 재고 비용 30% 절감", "연간 약 3,000만원 절감\n약 1년 만에 투자 회수", BLUE),
        ("⚡ 세금계산서 80% 개선", "월말 2일 → 2시간 완료\n정산 업무 부담 감소", PURPLE)
    ]

    y_pos = 2.3
    for title, detail, color in quantitative:
        add_text_box(slide, Inches(1.3), Inches(y_pos), Inches(5), Inches(0.3),
                     title, font_size=13, color=color, bold=True)
        add_text_box(slide, Inches(1.3), Inches(y_pos + 0.35), Inches(5), Inches(0.6),
                     detail, font_size=10, color=TEXT_SECONDARY)
        y_pos += 1.4

    # 우측: 정성적 효과
    add_bg_shape(slide, Inches(7.0), Inches(1.3), Inches(5.3), Inches(5.5), GRAY_50)
    add_text_box(slide, Inches(7.4), Inches(1.6), Inches(4.5), Inches(0.35),
                 "💡 정성적 효과", font_size=18, color=PURPLE, bold=True)

    qualitative = [
        ("🛡️ 재고 부족 사고 예방", "필수 의료 물품 결품 방지\n환자 진료 차질 최소화"),
        ("✅ 투명한 재무 관리", "거래 내역 완전 추적\n부정 사용 및 횡령 방지"),
        ("🤝 공급업체 협업 강화", "실시간 정보 공유\n장기 파트너십 구축")
    ]

    y_pos = 2.3
    for title, detail in qualitative:
        add_text_box(slide, Inches(7.5), Inches(y_pos), Inches(4.5), Inches(0.3),
                     title, font_size=13, color=PURPLE, bold=True)
        add_text_box(slide, Inches(7.5), Inches(y_pos + 0.35), Inches(4.5), Inches(0.6),
                     detail, font_size=10, color=TEXT_SECONDARY)
        y_pos += 1.4

    # ═══════════════════════════════════════════
    # SLIDE 8: Thank You
    # ═══════════════════════════════════════════
    print("슬라이드 08: Thank You")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_bg_shape(slide, Inches(0), Inches(0), SW, SH, WHITE)
    add_left_emerald_line(slide)

    # Thank You
    add_text_box(slide, Inches(0), Inches(2.2), SW, Inches(0.8),
                 "Thank You", font_size=56, color=EMERALD_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0), Inches(3.1), SW, Inches(0.4),
                 "정동병원의 디지털 전환을 위해 kitt AI가 함께하겠습니다", font_size=18, color=TEXT_SECONDARY,
                 alignment=PP_ALIGN.CENTER)

    # 구분선
    add_bg_shape(slide, Inches(5.5), Inches(3.8), Inches(2.3), Inches(0.03), BORDER_EMERALD)

    # 연락처
    add_text_box(slide, Inches(0), Inches(4.3), SW, Inches(0.35),
                 "kitt AI Intelligence", font_size=20, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)

    contact_lines = [
        "Email: inventory@kitt.ai",
        "Tel: 02-1234-5678",
        "Support: support@kitt.ai"
    ]
    for i, line in enumerate(contact_lines):
        add_text_box(slide, Inches(0), Inches(4.8 + i * 0.35), SW, Inches(0.3),
                     line, font_size=13, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0), Inches(6.3), SW, Inches(0.3),
                 "2026년 3월", font_size=12, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

    # ─── 저장 ───
    output_path = "/Users/jasonmac/kitt_agent/ai-agent/output/proposal/정동병원_재고관리_플랫폼_제안서.pptx"
    prs.save(output_path)
    print(f"\n✅ 프레젠테이션 생성 완료: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
